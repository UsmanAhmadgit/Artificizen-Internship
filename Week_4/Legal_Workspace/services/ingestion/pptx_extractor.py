import io
from pptx import Presentation
from services.ingestion.ocr_utils import extract_text_from_image
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_pptx_chunks(file_bytes: bytes) -> list[str]:
    if not file_bytes or len(file_bytes) == 0:
        return []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slide_chunks = []

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            slide_title = "Untitled Slide"

            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()

            content_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            content_parts.append(p_text)

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_text = extract_text_from_image(shape.image.blob)
                    if img_text:
                        clean_img_text = img_text.replace('\n', ' ')
                        content_parts.append(f"[Embedded Image Content]: {clean_img_text}")          

                elif shape.has_table:
                    table = shape.table
                    table_headers = []
                    for r_idx, row in enumerate(table.rows):
                        row_cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells if cell.text]
                        if not any(row_cells):
                            continue

                        if r_idx == 0:
                            table_headers = row_cells
                            content_parts.append("Table Header: " + " | ".join(table_headers))
                        else:
                            if table_headers and len(table_headers) == len(row_cells):
                                formatted_row = " | ".join([f"{h}: {v}" if h else v for h, v in zip(table_headers, row_cells)])
                            else:
                                formatted_row = " | ".join(row_cells)
                            content_parts.append("Table Row: " + formatted_row)

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            slide_body = "\n".join(content_parts)
            
            if slide_body or notes_text:
                full_slide_text = f"[Slide {slide_num}: {slide_title}]\n{slide_body}"
                if notes_text:
                    full_slide_text += f"\n[Slide {slide_num} Speaker Notes]: {notes_text}"

                if len(full_slide_text) > 800:
                    lines = full_slide_text.split("\n")
                    header_context = f"[Slide {slide_num}: {slide_title}]"
                    c_chunk = [header_context]
                    c_len = len(header_context)

                    for line in lines:
                        if line == header_context:
                            continue
                        if c_len + len(line) > 800 and len(c_chunk) > 1:
                            slide_chunks.append("\n".join(c_chunk))
                            c_chunk = [header_context, c_chunk[-1], line]
                            c_len = sum(len(x) for x in c_chunk)
                        else:
                            c_chunk.append(line)
                            c_len += len(line)
                    if len(c_chunk) > 1:
                        slide_chunks.append("\n".join(c_chunk))
                else:
                    slide_chunks.append(full_slide_text)

        return slide_chunks

    except Exception as e:
        raise RuntimeError(f"Failed to parse PowerPoint presentation: {str(e)}")